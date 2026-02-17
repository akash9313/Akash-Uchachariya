from pptx import Presentation

prs = Presentation()

def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content


# 1 Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Market Structures in Microeconomics"
slide.placeholders[1].text = "Perfect Competition, Monopoly, Monopolistic Competition & Oligopoly\nBased on Samuelson & Nordhaus"

# 2 Introduction
add_slide("Introduction to Market Structure",
"A market is a system where buyers and sellers interact to determine price and quantity through demand and supply. "
"Market structure refers to the degree of competition and monopoly power in a market and explains firm behaviour, pricing strategies and consumer welfare.")

# 3 Types of Market Structures
add_slide("Types of Market Structures",
"The four main forms of market structure are perfect competition, monopoly, monopolistic competition and oligopoly. "
"These structures differ in number of firms, product differentiation and control over price.")

# -------- PERFECT COMPETITION --------

# 4 Definition PC
add_slide("Perfect Competition: Definition",
"Perfect competition is a market where many buyers and sellers freely trade identical products and no single firm can influence the market price. "
"Example: Agricultural markets such as wheat, rice, fruits and vegetables.")

# 5 Features PC
add_slide("Features of Perfect Competition",
"Perfect competition is characterized by homogeneous products, a large number of buyers and sellers and free entry and exit. "
"Firms act as price takers, possess perfect knowledge and independently decide output to maximize profit.")

# 6 Industry Equilibrium PC
add_slide("Industry Equilibrium in Perfect Competition",
"Market equilibrium occurs where quantity demanded equals quantity supplied (Qd = Qs). "
"This determines equilibrium price and output, which becomes the market price accepted by all firms.")

# 7 Short Run Equilibrium PC
add_slide("Short-Run Firm Equilibrium in Perfect Competition",
"The firm produces where MR = MC and since the firm is a price taker, P = MR = AR. "
"At this output level, the firm maximizes profit or minimizes loss in the short run.")

# 8 Long Run Equilibrium PC
add_slide("Long-Run Firm Equilibrium in Perfect Competition",
"Entry and exit of firms continue until P = MR = MC = AC. "
"At this point firms earn only normal profit and the market reaches long-run equilibrium.")

# 9 Criticism PC
add_slide("Criticism / Limitations of Perfect Competition",
"Perfect competition is based on unrealistic assumptions such as perfect information, no advertising and no transport cost. "
"It is difficult to apply in real life and profits are not sustainable in the long run.")

# 10 Why Other Structures Exist
add_slide("Why Other Market Structures Exist",
"Perfect competition is rare because firms differentiate products and create entry barriers to gain price control. "
"This leads to the emergence of monopoly and monopolistic competition.")

# -------- MONOPOLY --------

# 11 Definition Monopoly
add_slide("Monopoly: Definition",
"A monopoly is a market structure where one firm is the sole producer of a good or service with no close substitutes and strong barriers to entry. "
"Example: Electricity supply, railways, patented medicines.")

# 12 Features Monopoly
add_slide("Features of Monopoly",
"Monopoly is characterized by a single seller, strong entry barriers, price making power, absence of close substitutes and a downward sloping demand curve.")

# 13 Short Run Monopoly
add_slide("Short-Run Monopoly Equilibrium",
"The monopolist maximizes profit where MR = MC and charges price from the demand curve, therefore P > MR. "
"This allows supernormal profit in the short run.")

# 14 Long Run Monopoly
add_slide("Long-Run Monopoly Equilibrium",
"In the long run entry is blocked, so MR = MC and P > AC. "
"The firm continues earning abnormal profit even in the long run.")

# 15 Criticism Monopoly
add_slide("Criticism of Monopoly",
"Monopoly leads to higher prices, price discrimination and reduced consumer welfare. "
"Lack of competition may reduce quality and encourage unfair trade practices.")

# 16 Classification Monopoly
add_slide("Classification of Monopoly",
"Monopoly may be natural monopoly, legal monopoly or coercive monopoly depending on the source of market power and barriers to entry.")

# -------- MONOPOLISTIC COMPETITION --------

# 17 Definition MC
add_slide("Monopolistic Competition: Definition",
"Monopolistic competition is a market structure between perfect competition and monopoly where firms sell differentiated products and have some price control. "
"Example: Restaurants, clothing brands, salons.")

# 18 Features MC
add_slide("Features of Monopolistic Competition",
"It is characterized by product differentiation, many small firms, low entry barriers, selling costs and a downward sloping demand curve.")

# 19 Short Run MC
add_slide("Short-Run Equilibrium in Monopolistic Competition",
"Each firm maximizes profit where MR = MC and because the demand curve slopes downward, P > MR allowing supernormal profit in the short run.")

# 20 Long Run MC
add_slide("Long-Run Equilibrium in Monopolistic Competition",
"In the long run new firms enter reducing profits until MR = MC and P = AC. "
"Firms earn only normal profit.")

# 21 Criticism MC
add_slide("Criticism of Monopolistic Competition",
"This market structure leads to higher prices due to branding, excess advertising costs and consumer confusion caused by product differentiation.")

# -------- OLIGOPOLY --------

# 22 Definition Oligopoly
add_slide("Oligopoly: Definition",
"Oligopoly is a market structure where a small number of firms dominate the market and each firm’s decisions affect the others. "
"Example: Telecom companies, automobile industry, airlines.")

# 23 Features Oligopoly
add_slide("Features of Oligopoly",
"Oligopoly is characterized by interdependence, few dominant firms, heavy advertising, high barriers to entry and strategic decision making.")

# 24 Kinked Demand Curve
add_slide("Kinked Demand Curve Model",
"Firms believe competitors will follow price decreases but ignore price increases, creating a kinked demand curve and price rigidity. "
"Equilibrium condition remains MR = MC.")

# 25 Types of Oligopoly
add_slide("Types of Oligopoly",
"Oligopoly may be pure oligopoly, differentiated oligopoly, collusive oligopoly or non-collusive oligopoly depending on product nature and firm behaviour.")

# 26 Collusion & Cartel
add_slide("Collusion, Cartel and Price Leadership",
"Firms may collude to fix prices, form cartels to behave like a monopoly or follow price leadership where one dominant firm sets the price.")

# 27 Criticism Oligopoly
add_slide("Criticism of Oligopoly",
"Oligopoly reduces competition and choice, may result in higher prices and can slow innovation due to reduced competitive pressure.")

# -------- CONCLUSION --------

# 28 Comparison
add_slide("Comparative Overview",
"Perfect competition represents maximum competition, monopoly shows absence of competition, while monopolistic competition and oligopoly lie between these extremes.")

# 29 Importance
add_slide("Importance of Market Structures",
"Understanding market structures helps analyse firm behaviour, pricing decisions and their impact on consumers and economic efficiency.")

# 30 Real World Relevance
add_slide("Real World Relevance",
"Most real markets resemble monopolistic competition or oligopoly rather than perfect competition or pure monopoly.")

# 31 Key Takeaways
add_slide("Key Takeaways",
"Different market structures influence price, output and profits differently based on competition level and market power.")

# 32 Final Conclusion
add_slide("Conclusion",
"Market structures explain how competition affects firm behaviour, pricing and profits. "
"They help understand decision making, price determination and the impact on consumers and economic efficiency.")

prs.save("Market_Structure_32_Slides_PDF_Style.pptx")
print("PPT with PDF-style structure created successfully!")
